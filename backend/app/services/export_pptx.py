"""Build a PowerPoint (.pptx) presentation from OKR objectives."""
from bson import ObjectId
from pptx import Presentation
from pptx.util import Inches, Pt


def _parse_oid(value):
    try:
        return ObjectId(value)
    except Exception:
        return None


def build_okr_pptx(db, objective_ids, output_stream, narrative=None):
    """
    Write a PowerPoint presentation to output_stream (BytesIO).
    objective_ids: list of objective ID strings (from request).
    narrative: optional AI-generated script text; when set, add a "Presentation script" slide after the title.
    Creates: title slide + [narrative slide if narrative] + one slide per objective (title, status, key results with scores).
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Load objectives and key results
    items = []
    for oid_str in objective_ids:
        oid = _parse_oid(oid_str)
        if not oid:
            continue
        obj = db.objectives.find_one({'_id': oid})
        if obj:
            krs = list(db.key_results.find({'objectiveId': oid}))
            items.append((obj, krs))

    if not items:
        # Empty presentation: one slide with title
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        except IndexError:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        box.text_frame.paragraphs[0].text = "OKR Presentation"
        box.text_frame.paragraphs[0].font.size = Pt(44)
        box.text_frame.paragraphs[0].font.bold = True
        prs.save(output_stream)
        return

    # Title slide
    try:
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "OKR Presentation"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"{len(items)} objective(s)"
    except (IndexError, AttributeError, KeyError):
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        box.text_frame.paragraphs[0].text = "OKR Presentation"
        box.text_frame.paragraphs[0].font.size = Pt(44)

    # Optional narrative/script slide (after title, before objectives)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    if narrative and isinstance(narrative, str) and narrative.strip():
        script_slide = prs.slides.add_slide(blank_layout)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5)
        title_box = script_slide.shapes.add_textbox(left, top, width, height)
        title_box.text_frame.paragraphs[0].text = "Presentation script"
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True
        top += Inches(0.7)
        body_box = script_slide.shapes.add_textbox(left, top, width, Inches(6))
        tf = body_box.text_frame
        tf.word_wrap = True
        lines = narrative.strip().split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line[:500]
            p.font.size = Pt(14)
            p.space_after = Pt(6)

    # One slide per objective
    for obj, krs in items:
        slide = prs.slides.add_slide(blank_layout)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = (obj.get('title') or 'Untitled')[:200]
        p.font.size = Pt(28)
        p.font.bold = True

        status = obj.get('status') or 'draft'
        top += Inches(1.2)
        sub = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        sub.text_frame.paragraphs[0].text = f"Status: {status}"
        sub.text_frame.paragraphs[0].font.size = Pt(14)
        top += Inches(0.6)

        for kr in krs[:10]:
            score = kr.get('score')
            score_str = f" {int((score or 0) * 100)}%" if score is not None else ""
            line = f"• {(kr.get('title') or '')[:100]}{score_str}"
            box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            box.text_frame.paragraphs[0].text = line
            box.text_frame.paragraphs[0].font.size = Pt(14)
            top += Inches(0.45)

    prs.save(output_stream)
